import os
import sys
import json
import logging
import asyncio
import traceback
from pathlib import Path
from logging.handlers import TimedRotatingFileHandler
from pptx import Presentation
import google.generativeai as genai
from dotenv import load_dotenv

BASE_DIR = Path(__file__).parent.parent.absolute()
sys.path.insert(0, str(BASE_DIR))

from database import get_session, Upload, UploadStatus


# ---------------------------------------------------
# Directory setup
# ---------------------------------------------------
UPLOADS_DIR = os.environ.get("UPLOADS_DIR", os.path.join(BASE_DIR, "shared", "uploads"))
OUTPUTS_DIR = os.environ.get("OUTPUTS_DIR", os.path.join(BASE_DIR, "shared", "outputs"))
LOGS_DIR = os.environ.get("LOGS_DIR", os.path.join(BASE_DIR, "shared", "logs", "explainer"))
ORIGINAL_PROJECT_DIR = os.path.join(BASE_DIR, "gemini-explainer")

print(f"DEBUG: BASE_DIR = {BASE_DIR}")
print(f"DEBUG: UPLOADS_DIR = {UPLOADS_DIR}")

for directory in [UPLOADS_DIR, OUTPUTS_DIR, LOGS_DIR]:
    os.makedirs(directory, exist_ok=True)

# Load .env
dotenv_path = os.path.join(ORIGINAL_PROJECT_DIR, ".env")
load_dotenv(dotenv_path)


# ---------------------------------------------------
# SUMMARY LEVEL PROMPTS & LANGUAGE TEMPLATES
# ---------------------------------------------------
PROMPT_TEMPLATES = {
    "beginner": {
        "en": "Explain this PowerPoint slide in simple, easy-to-understand terms. Avoid technical jargon and use everyday language that anyone can follow.",
        "he": "הסבר את שקף הPowerPoint הזה במילים פשוטות וקלות להבנה. הימנע מטרמינולוגיה טכנית והשתמש בשפה יומיומית.",
        "ru": "Объясните этот слайд PowerPoint простыми и понятными словами. Избегайте технической терминологии и используйте повседневный язык.",
        "es": "Explique esta diapositiva de PowerPoint con palabras simples y fáciles de entender. Evite la jerga técnica y use lenguaje cotidiano.",
    },
    "comprehensive": {
        "en": "Provide a detailed and thorough explanation of this PowerPoint slide. Include all key concepts, important details, relevant context, and any notable implications.",
        "he": "תן הסבר מפורט ויסודי של שקף PowerPoint זה. כלול את כל המושגים העיקריים, פרטים חשובים, הקשר רלוונטי וכל השלכה ניכרת.",
        "ru": "Предоставьте подробное и тщательное объяснение этого слайда PowerPoint. Включите все ключевые концепции, важные детали, актуальный контекст и любые значительные последствия.",
        "es": "Proporcione una explicación detallada y exhaustiva de esta diapositiva de PowerPoint. Incluya todos los conceptos clave, detalles importantes, contexto relevante e implicaciones notables.",
    },
    "executive": {
        "en": "Provide a concise executive summary of this PowerPoint slide in 2-3 sentences. Focus only on the most critical information and main takeaways.",
        "he": "תן סיכום ביצועי תמציתי של שקף PowerPoint זה ב-2-3 משפטים. התמקד רק במידע הקריטי ביותר והנקודות העיקריות.",
        "ru": "Предоставьте краткий исполнительный резюме этого слайда PowerPoint в 2-3 предложениях. Сосредоточьтесь только на наиболее важной информации и основных выводах.",
        "es": "Proporcione un resumen ejecutivo conciso de esta diapositiva de PowerPoint en 2-3 oraciones. Enfóquese solo en la información más crítica y los puntos clave.",
    },
}


def get_prompt_for_level_and_language(slide_text, summary_level, language):
    """Get the appropriate prompt based on summary level and language"""
    base_prompt = PROMPT_TEMPLATES.get(summary_level, PROMPT_TEMPLATES["comprehensive"]).get(language, PROMPT_TEMPLATES["comprehensive"]["en"])
    return f"{base_prompt}\n\nSlide content: {slide_text}"


# ---------------------------------------------------
# TEXT EXTRACTION
# ---------------------------------------------------
async def extract_text_from_slide(slide):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + "\n"

    text = " ".join(text.split())
    return text.strip()


async def parse_presentation(pptx_path):
    presentation = Presentation(pptx_path)
    slides_with_text = []

    for idx, slide in enumerate(presentation.slides):
        text = await extract_text_from_slide(slide)
        if text:
            slides_with_text.append({
                "slide_number": idx + 1,
                "text": text
            })

    return slides_with_text


# ---------------------------------------------------
# GEMINI API CALL – FIXED VERSION
# ---------------------------------------------------
import google.generativeai as genai

async def get_explanation(slide_text, summary_level, language, timeout=30):
    try:
        prompt = get_prompt_for_level_and_language(slide_text, summary_level, language)

        # הגדרת ה־API Key
        genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

        # יצירת מודל
        model = genai.GenerativeModel("gemini-2.0-flash-lite-preview-02-05")

        def call_gemini():
            response = model.generate_content(prompt)
            return response.text

        explanation = await asyncio.wait_for(
            asyncio.to_thread(call_gemini),
            timeout=timeout
        )

        return explanation

    except asyncio.TimeoutError:
        return "Error: Gemini API request timed out."
    except Exception as e:
        return f"Error explaining this slide: {str(e)}"



# ---------------------------------------------------
# SLIDE PROCESSING
# ---------------------------------------------------
async def process_slides(slides, summary_level, language, delay_seconds=3, max_parallel=3):
    results = []

    for i in range(0, len(slides), max_parallel):
        batch = slides[i:i + max_parallel]
        tasks = []

        for slide in batch:
            task = asyncio.create_task(get_explanation(slide["text"], summary_level, language))
            tasks.append((slide["slide_number"], task))

        # Collect results
        for slide_number, task in tasks:
            try:
                explanation = await task
                results.append({
                    "slide_number": slide_number,
                    "explanation": explanation
                })
            except Exception as e:
                results.append({
                    "slide_number": slide_number,
                    "explanation": f"Error processing slide {slide_number}: {str(e)}"
                })

        if i + max_parallel < len(slides):
            logger.info(f"Waiting {delay_seconds} seconds before next batch...")
            await asyncio.sleep(delay_seconds)

    results.sort(key=lambda x: x["slide_number"])
    return results


# ---------------------------------------------------
# LOGGER
# ---------------------------------------------------
def setup_logger():
    logger = logging.getLogger('explainer')
    logger.setLevel(logging.INFO)

    console = logging.StreamHandler()
    console.setFormatter(logging.Formatter(
        '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    ))
    logger.addHandler(console)

    file_handler = TimedRotatingFileHandler(
        os.path.join(LOGS_DIR, "explainer.log"),
        when="midnight",
        backupCount=5
    )
    file_handler.setFormatter(logging.Formatter(
        '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    ))
    logger.addHandler(file_handler)

    return logger


logger = setup_logger()


# ---------------------------------------------------
# Upload Processing
# ---------------------------------------------------
SLEEP_INTERVAL = 10


async def process_upload(upload):
    logger.info(f"Processing upload: {upload.uid} ({upload.filename}) - Level: {upload.summary_level}")

    session = get_session()
    try:
        upload_db = session.query(Upload).filter(Upload.id == upload.id).first()
        upload_db.set_status_processing()
        session.commit()

        if not os.path.exists(upload.upload_path):
            msg = f"Upload file not found: {upload.upload_path}"
            logger.error(msg)
            upload_db.set_status_failed(msg)
            session.commit()
            return

        slides = await parse_presentation(upload.upload_path)

        if not slides:
            msg = f"No slides with text in {upload.filename}"
            logger.warning(msg)
            upload_db.set_status_failed(msg)
            session.commit()
            return

        logger.info(f"Found {len(slides)} slides, processing with level: {upload.summary_level} in language: {upload.language}…")

        explanations = await process_slides(slides, upload.summary_level, upload.language)

        with open(upload.output_path, "w", encoding="utf-8") as f:
            json.dump(explanations, f, indent=2, ensure_ascii=False)

        upload_db.set_status_completed()
        session.commit()

        logger.info(f"Successfully processed {upload.uid}")

    except Exception as e:
        msg = f"Error processing upload {upload.uid}: {str(e)}"
        logger.error(msg)
        logger.error(traceback.format_exc())

        upload_db = session.query(Upload).filter(Upload.id == upload.id).first()
        upload_db.set_status_failed(msg)
        session.commit()

    finally:
        session.close()


async def find_pending_uploads():
    session = get_session()
    try:
        return session.query(Upload) \
            .filter(Upload.status == UploadStatus.UPLOADED) \
            .order_by(Upload.upload_time.asc()) \
            .all()
    except Exception as e:
        logger.error(f"Error finding uploads: {str(e)}")
        return []
    finally:
        session.close()


# ---------------------------------------------------
# MAIN LOOP
# ---------------------------------------------------
async def main_loop():
    logger.info("Explainer service started")

    while True:
        try:
            pending = await find_pending_uploads()

            if pending:
                logger.info(f"Found {len(pending)} pending uploads")
                for upload in pending:
                    await process_upload(upload)

            await asyncio.sleep(SLEEP_INTERVAL)

        except Exception as e:
            logger.error(f"Error in main loop: {str(e)}")
            logger.error(traceback.format_exc())
            await asyncio.sleep(SLEEP_INTERVAL)



if __name__ == "__main__":
    asyncio.run(main_loop())