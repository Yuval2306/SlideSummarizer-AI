from pptx import Presentation

async def extract_text_from_slide(slide):
    """Extract all text from a PowerPoint slide."""
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + "\n"

    # Clean up whitespace
    text = " ".join(text.split())
    return text.strip()


async def parse_presentation(pptx_path):
    """Parse a PowerPoint presentation and extract text from each slide."""
    presentation = Presentation(pptx_path)
    slides_with_text = []

    for idx, slide in enumerate(presentation.slides):
        text = await extract_text_from_slide(slide)
        if text:
            slides_with_text.append({
                "slide_number": idx + 1,
                "text": text
            })

    return slides_with_text


